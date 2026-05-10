"""Generate a merged PPT (no duplicate sections): all figures + concise explanations + conclusion."""

from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


REPO_ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = REPO_ROOT / "figures"
OUT_PPT = REPO_ROOT / "slides" / "term_project_presentation_SYNGAS--group 6 Ying.pptx"
RESULTS_PATH = REPO_ROOT / "results_summary.txt"

# K-State inspired palette
KSTATE_PURPLE = RGBColor(81, 40, 136)
KSTATE_LAVENDER = RGBColor(232, 226, 240)
TEXT_DARK = RGBColor(37, 31, 45)
WHITE = RGBColor(255, 255, 255)

FIG_DATA = [
    (
        "fig1_data_distributions.png",
        "Figure 1: Data Distributions",
        [
            "Shows spread/skewness of all features and target.",
            "Confirms heterogeneous feature scales.",
            "Supports preprocessing and scaling decisions.",
        ],
    ),
    (
        "fig2_feature_vs_target.png",
        "Figure 2: Feature vs Target",
        [
            "Visual evidence of nonlinear relationships.",
            "Motivates RF/MLP beyond linear baseline.",
            "Helps identify interaction effects.",
        ],
    ),
    (
        "fig3_mlp_training_loss.png",
        "Figure 3: MLP Training Loss Curve",
        [
            "Tracks optimization behavior across epochs.",
            "Shows convergence trend and sensitivity.",
            "Used for training-dynamics discussion.",
        ],
    ),
    (
        "fig4_train_vs_test_performance.png",
        "Figure 4: Train vs Test Performance",
        [
            "Compares generalization across LR/RF/MLP.",
            "RF has strongest test metrics.",
            "Useful for overfit/underfit diagnosis.",
        ],
    ),
    (
        "fig5_scatter_LR_train_test.png",
        "Figure 5: LR Actual vs Predicted",
        [
            "Baseline behavior on train/test splits.",
            "Larger deviation from diagonal than RF/MLP.",
            "Reference for significance tests.",
        ],
    ),
    (
        "fig6_scatter_RF_train_test.png",
        "Figure 6: RF Actual vs Predicted",
        [
            "Predictions cluster close to diagonal.",
            "Best test performance among compared models.",
            "Primary deployment candidate evidence.",
        ],
    ),
    (
        "fig7_scatter_MLP_train_test.png",
        "Figure 7: MLP Actual vs Predicted",
        [
            "Nonlinear fit is competitive with RF.",
            "Sensitivity appears higher than RF.",
            "Interpreted with scaling ablation results.",
        ],
    ),
    (
        "fig8_RF_feature_importance.png",
        "Figure 8: RF Feature Importance",
        [
            "Ranks feature contributions in RF.",
            "Improves interpretability for operations.",
            "Guides future feature engineering.",
        ],
    ),
    (
        "fig9_error_comparison.png",
        "Figure 9: Error Comparison",
        [
            "Compares test-set error distributions.",
            "Paired sample view supports hypothesis testing.",
            "RF errors are generally lower than LR.",
        ],
    ),
    (
        "fig10_robustness_boxplot.png",
        "Figure 10: Robustness Boxplot",
        [
            "RMSE distribution over 30 random splits.",
            "Measures stability beyond one split.",
            "RF remains most stable overall.",
        ],
    ),
    (
        "fig11_robustness_line.png",
        "Figure 11: Robustness Across Seeds",
        [
            "Seed-level RMSE trajectories by model.",
            "Shows variance behavior directly.",
            "Complements Figure 10 statistics.",
        ],
    ),
    (
        "fig12_scaling_ablation_rmse.png",
        "Figure 12: Scaling Ablation (4 Combinations)",
        [
            "Tests all required scaler on/off combinations.",
            "RF unscaled performs best in this dataset.",
            "Documents preprocessing impact explicitly.",
        ],
    ),
    (
        "fig13_runtime_ablation.png",
        "Figure 13: Runtime Analysis",
        [
            "Wall-clock training and inference comparison.",
            "Inference is sub-millisecond per sample.",
            "Runtime is acceptable for deployment.",
        ],
    ),
]


def parse_summary(path: Path) -> dict:
    text = path.read_text(encoding="utf-8")
    out = {}
    pats = {
        "rf": r"Random Forest:\s*\n\s*Train:.*\n\s*Test:\s*RMSE=([0-9.]+),\s*MAE=([0-9.]+),\s*R²=([0-9.]+)",
        "lr": r"Linear Regression:\s*\n\s*Train:.*\n\s*Test:\s*RMSE=([0-9.]+),\s*MAE=([0-9.]+),\s*R²=([0-9.]+)",
        "mlp": r"MLP:\s*\n\s*Train:.*\n\s*Test:\s*RMSE=([0-9.]+),\s*MAE=([0-9.]+),\s*R²=([0-9.]+)",
        "t": r"Paired t-test:\s*t=([0-9.\-]+),\s*p=([0-9.eE\-+]+)",
        "w": r"Wilcoxon:\s*W=([0-9.\-]+),\s*p=([0-9.eE\-+]+)",
    }
    for k, p in pats.items():
        m = re.search(p, text)
        if m:
            out[k] = m.groups()
    return out


def add_title(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_theme_frame(s, "Title")

    hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(9.0), Inches(4.9))
    hero.fill.solid()
    hero.fill.fore_color.rgb = KSTATE_PURPLE
    hero.line.fill.background()

    tbox = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(2.1))
    tf = tbox.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = "Ethanol Concentration Prediction\nin Syngas Fermentation"
    p0.font.size = Pt(32)
    p0.font.bold = True
    p0.font.color.rgb = WHITE
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "CIS 732/830 Final Project"
    p1.font.size = Pt(18)
    p1.font.color.rgb = KSTATE_LAVENDER
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Ying Tan | K-State Theme Edition"
    p2.font.size = Pt(14)
    p2.font.color.rgb = KSTATE_LAVENDER
    p2.alignment = PP_ALIGN.CENTER


def add_overview(prs: Presentation, sm: dict):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_theme_frame(s, "Overview")

    tf = add_title_box(s, "Executive Summary")

    body_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.45), Inches(8.8), Inches(4.9))
    body_box.fill.solid()
    body_box.fill.fore_color.rgb = RGBColor(246, 244, 250)
    body_box.line.color.rgb = KSTATE_LAVENDER

    body_tf = s.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(8.2), Inches(4.2)).text_frame
    body_tf.clear()
    tf.clear()

    rf = sm.get("rf", ("?", "?", "?"))
    lr = sm.get("lr", ("?", "?", "?"))
    mlp = sm.get("mlp", ("?", "?", "?"))
    t = sm.get("t", ("?", "?"))
    w = sm.get("w", ("?", "?"))

    lines = [
        "Dataset: 176 samples, 7 features, target = ethanol (mM)",
        f"RF test: RMSE={rf[0]}, MAE={rf[1]}, R2={rf[2]} (best)",
        f"MLP test: RMSE={mlp[0]}, MAE={mlp[1]}, R2={mlp[2]}",
        f"LR test: RMSE={lr[0]}, MAE={lr[1]}, R2={lr[2]}",
        f"Statistical tests: paired t-test p={t[1]}, Wilcoxon p={w[1]}",
        "Slides 3-15 cover Figure 1 to Figure 13 with explanations.",
    ]
    for i, line in enumerate(lines):
        p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18 if i == 0 else 15)
        p.font.color.rgb = TEXT_DARK
        if i > 0:
            p.level = 0


def add_figure_slide(prs: Presentation, fig: str, title: str, bullets: list[str]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_theme_frame(s, "Figure")

    add_title_box(s, title)

    # Left panel for image
    img_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(1.15), Inches(6.55), Inches(5.35))
    img_card.fill.solid()
    img_card.fill.fore_color.rgb = WHITE
    img_card.line.color.rgb = KSTATE_LAVENDER

    s.shapes.add_picture(str(FIG_DIR / fig), Inches(0.52), Inches(1.38), width=Inches(6.2), height=Inches(4.9))

    # Right panel for explanation text
    note_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(1.15), Inches(2.6), Inches(5.35))
    note_card.fill.solid()
    note_card.fill.fore_color.rgb = RGBColor(247, 244, 252)
    note_card.line.color.rgb = KSTATE_LAVENDER

    nbox = s.shapes.add_textbox(Inches(7.22), Inches(1.35), Inches(2.28), Inches(4.95))
    ntf = nbox.text_frame
    ntf.word_wrap = True
    ntf.clear()
    for i, b in enumerate(bullets):
        p = ntf.paragraphs[0] if i == 0 else ntf.add_paragraph()
        p.text = f"- {b}"
        p.font.size = Pt(13 if i == 0 else 12)
        p.font.color.rgb = TEXT_DARK

    cap = s.shapes.add_textbox(Inches(0.45), Inches(6.15), Inches(6.4), Inches(0.28)).text_frame
    cap.text = "All values and figures are generated from scripts/syngas_analysis.py"
    cap.paragraphs[0].font.size = Pt(10)
    cap.paragraphs[0].font.color.rgb = RGBColor(98, 93, 108)


def add_conclusion(prs: Presentation, sm: dict):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_theme_frame(s, "Conclusion")
    tf = add_title_box(s, "Conclusion and Recommendation")

    body_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.45), Inches(8.8), Inches(4.9))
    body_box.fill.solid()
    body_box.fill.fore_color.rgb = RGBColor(246, 244, 250)
    body_box.line.color.rgb = KSTATE_LAVENDER

    tf = s.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(8.2), Inches(4.2)).text_frame
    tf.clear()

    rf = sm.get("rf", ("?", "?", "?"))
    t = sm.get("t", ("?", "?"))
    w = sm.get("w", ("?", "?"))
    lines = [
        "Random Forest is selected as the primary model for this dataset.",
        f"Best fixed-split result: RF RMSE={rf[0]}, R2={rf[2]}",
        "30-split robustness supports stable generalization.",
        "Scaling ablation and runtime analyses were fully completed.",
        f"Significance confirmed: t-test p={t[1]}, Wilcoxon p={w[1]}",
        "Final recommendation: deploy RF first, keep MLP for future extensions.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18 if i == 0 else 15)
        p.font.color.rgb = TEXT_DARK


def add_theme_frame(slide, tag: str):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10.0), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(252, 251, 255)
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10.0), Inches(0.45))
    top.fill.solid()
    top.fill.fore_color.rgb = KSTATE_PURPLE
    top.line.fill.background()

    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(10.0), Inches(0.3))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(236, 230, 245)
    footer.line.fill.background()

    footer_tf = slide.shapes.add_textbox(Inches(0.28), Inches(7.22), Inches(9.4), Inches(0.22)).text_frame
    footer_tf.text = f"Kansas State University | CIS 732/830 | {tag}"
    footer_tf.paragraphs[0].font.size = Pt(10)
    footer_tf.paragraphs[0].font.color.rgb = RGBColor(84, 76, 102)


def add_title_box(slide, text: str):
    tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.52), Inches(9.2), Inches(0.55))
    tf = tbox.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = KSTATE_PURPLE
    return tf


def main():
    missing = [f for f, _, _ in FIG_DATA if not (FIG_DIR / f).exists()]
    if missing:
        raise FileNotFoundError(f"Missing figures: {missing}")

    prs = Presentation()
    sm = parse_summary(RESULTS_PATH)
    add_title(prs)
    add_overview(prs, sm)
    for fig, title, bullets in FIG_DATA:
        add_figure_slide(prs, fig, title, bullets)
    add_conclusion(prs, sm)

    prs.save(OUT_PPT)
    print("Saved:", OUT_PPT)
    print("Slides:", len(prs.slides))


if __name__ == "__main__":
    main()
